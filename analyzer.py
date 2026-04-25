import os
import nltk
from pptx import Presentation
from textblob import TextBlob

# Ensure NLTK resources are downloaded
try:
    nltk.data.find('tokenizers/punkt')
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('punkt', quiet=True)
    nltk.download('stopwords', quiet=True)

from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize

def extract_text_from_ppt(filepath):
    """
    Extracts text from all slides in a .pptx file.
    Returns a dictionary mapping slide_number -> text
    """
    prs = Presentation(filepath)
    slides_data = {}
    
    for i, slide in enumerate(prs.slides):
        text_runs = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
        
        # Join text and clean up excessive newlines
        slide_text = " ".join(text_runs).strip()
        slides_data[i + 1] = slide_text
        
    return slides_data

def fix_presentation_and_save(input_filepath, output_filepath):
    """
    Creates a new PPT with corrected text while keeping the same design.
    """
    try:
        prs = Presentation(input_filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text and run.text.strip():
                                # Correct the text using TextBlob
                                corrected_text = str(TextBlob(run.text).correct())
                                run.text = corrected_text
        prs.save(output_filepath)
        return True
    except Exception as e:
        print(f"Error fixing presentation: {e}")
        return False

def analyze_slide_length(text):
    """
    Analyzes if the slide has too much or too little text.
    Ideal length: between 10 and 60 words for a good presentation.
    """
    words = word_tokenize(text)
    word_count = len(words)
    
    if word_count == 0:
        return "Empty slide.", 0
    elif word_count < 10:
        return "Too little text. Consider adding more context.", word_count
    elif word_count > 60:
        return "Too much text. Slide might look cluttered. Consider using bullet points.", word_count
    else:
        return "Good content length.", word_count

def analyze_grammar(text):
    """
    Checks for spelling mistakes using TextBlob for an instant, lightweight check.
    """
    if not text.strip():
        return [], 0
    
    blob = TextBlob(text)
    errors = []
    
    # We check for words that TextBlob thinks are misspelled
    for word in blob.words:
        # A simple check: if the word's spellcheck doesn't match the original, it might be misspelled.
        # We only take the top suggestion.
        corrected_word = word.correct()
        if word.lower() != corrected_word.lower() and len(word) > 2:
            errors.append({
                'message': f"Possible misspelling: '{word}'",
                'context': f"...{word}...",
                'replacements': [str(corrected_word)]
            })
            
    # Textblob may find many false positives, so let's limit errors per slide
    return errors[:5], min(len(errors), 5)

def analyze_flow(slides_data):
    """
    Basic flow analysis: checks if consecutive slides share common keywords
    to ensure smooth transition.
    """
    stop_words = set(stopwords.words('english'))
    flow_feedback = []
    
    prev_keywords = set()
    
    for slide_num, text in slides_data.items():
        if not text.strip():
            flow_feedback.append(f"Slide {slide_num} is empty, breaking the flow.")
            prev_keywords = set()
            continue
            
        words = word_tokenize(text.lower())
        keywords = {w for w in words if w.isalnum() and w not in stop_words}
        
        if slide_num > 1 and prev_keywords:
            overlap = prev_keywords.intersection(keywords)
            if not overlap:
                flow_feedback.append(f"Weak transition from Slide {slide_num-1} to Slide {slide_num}. Consider adding linking concepts.")
                
        prev_keywords = keywords
        
    if not flow_feedback:
        flow_feedback.append("Good flow and structure detected overall.")
        
    return flow_feedback

def get_top_keywords(text, n=6):
    stop_words = set(stopwords.words('english'))
    words = word_tokenize(text.lower())
    keywords = [w for w in words if w.isalpha() and w not in stop_words and len(w) > 3]
    freq = {}
    for w in keywords:
        freq[w] = freq.get(w, 0) + 1
    sorted_kw = sorted(freq.items(), key=lambda x: x[1], reverse=True)
    return [{'word': w, 'count': c} for w, c in sorted_kw[:n]]

def get_readability_score(text):
    words = word_tokenize(text)
    sentences = text.split('.')
    sentences = [s.strip() for s in sentences if s.strip()]
    word_count = len(words)
    sentence_count = max(len(sentences), 1)
    syllable_count = sum(count_syllables(w) for w in words)
    if word_count == 0:
        return 0, 'N/A'
    score = 206.835 - 1.015 * (word_count / sentence_count) - 84.6 * (syllable_count / max(word_count, 1))
    score = max(0, min(100, round(score, 1)))
    if score >= 70: label = 'Easy'
    elif score >= 50: label = 'Moderate'
    else: label = 'Difficult'
    return score, label

def count_syllables(word):
    word = word.lower()
    vowels = 'aeiouy'
    count = 0
    prev_vowel = False
    for ch in word:
        is_vowel = ch in vowels
        if is_vowel and not prev_vowel:
            count += 1
        prev_vowel = is_vowel
    return max(1, count)

def analyze_presentation(filepath, fixed_filepath=None):
    slides_data = extract_text_from_ppt(filepath)
    analysis_results = {
        'slide_analysis': [],
        'overall_score': 10.0,
        'flow_feedback': [],
        'total_grammar_errors': 0
    }
    total_slides = len(slides_data)
    if total_slides == 0:
        return {"error": "No slides found in the presentation."}
    score_deductions = 0
    for slide_num, text in slides_data.items():
        length_feedback, word_count = analyze_slide_length(text)
        grammar_errors, error_count = analyze_grammar(text)
        keywords = get_top_keywords(text)
        readability_score, readability_label = get_readability_score(text)
        analysis_results['total_grammar_errors'] += error_count
        if "Too little text" in length_feedback or "Too much text" in length_feedback:
            score_deductions += 0.5
        score_deductions += min(error_count * 0.2, 1.0)
        analysis_results['slide_analysis'].append({
            'slide_number': slide_num,
            'text_snippet': text[:100] + '...' if len(text) > 100 else text,
            'word_count': word_count,
            'length_feedback': length_feedback,
            'grammar_errors': grammar_errors,
            'error_count': error_count,
            'keywords': keywords,
            'readability_score': readability_score,
            'readability_label': readability_label
        })
    flow_feedback = analyze_flow(slides_data)
    analysis_results['flow_feedback'] = flow_feedback
    if len(flow_feedback) > 1 or "Weak transition" in str(flow_feedback):
        score_deductions += len([f for f in flow_feedback if "Weak transition" in f]) * 0.5
    final_score = max(0.0, 10.0 - (score_deductions / max(1, (total_slides / 5))))
    analysis_results['overall_score'] = round(final_score, 1)
    
    # Generate the fixed presentation
    if fixed_filepath:
        success = fix_presentation_and_save(filepath, fixed_filepath)
        if success:
            analysis_results['fixed_file_available'] = True
        else:
            analysis_results['fixed_file_available'] = False

    return analysis_results
